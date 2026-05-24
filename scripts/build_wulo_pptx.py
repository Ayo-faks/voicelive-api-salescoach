"""Build a clickable PPTX from the Wulo HTML deck.

Each slide becomes a 1920x1080 image with transparent rectangles overlaying
every <a href="http..."> so the links remain clickable in PowerPoint.
"""
from __future__ import annotations

from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "docs" / "wulo-deck"
INDEX = DECK / "index.html"
EXPORTS = DECK / "exports" / "1920x1080"
OUT = DECK / "exports" / "wulo-deck.pptx"

VIEW_W, VIEW_H = 1920, 1080
SLIDE_W = Emu(12192000)  # 13.333"
SLIDE_H = Emu(6858000)   # 7.5"
PX_TO_EMU = SLIDE_W / VIEW_W


def collect() -> list[tuple[Path, list[dict]]]:
    out: list[tuple[Path, list[dict]]] = []
    with sync_playwright() as p:
        browser = p.chromium.launch()
        ctx = browser.new_context(viewport={"width": VIEW_W, "height": VIEW_H}, device_scale_factor=2)
        page = ctx.new_page()
        page.goto(INDEX.as_uri())
        page.evaluate("document.body.classList.add('export')")
        page.wait_for_selector(".slide")
        count = page.evaluate("document.querySelectorAll('.slide').length")
        for i in range(count):
            meta = page.evaluate(
                """(i) => {
                    const ss = document.querySelectorAll('.slide');
                    ss.forEach(s => s.classList.remove('active'));
                    const s = ss[i];
                    s.classList.add('active');
                    return { slide: s.dataset.slide, name: s.dataset.name };
                }""",
                i,
            )
            page.wait_for_timeout(120)
            links = page.evaluate(
                """() => {
                    const active = document.querySelector('.slide.active');
                    const out = [];
                    active.querySelectorAll('a[href^="http"]').forEach(a => {
                        const r = a.getBoundingClientRect();
                        if (r.width > 0 && r.height > 0) {
                            out.push({ href: a.href, x: r.left, y: r.top, w: r.width, h: r.height });
                        }
                    });
                    return out;
                }"""
            )
            img = EXPORTS / f"{meta['slide']}-{meta['name']}.png"
            out.append((img, links))
        browser.close()
    return out


def build(items: list[tuple[Path, list[dict]]]) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for img_path, links in items:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(img_path), 0, 0, width=SLIDE_W, height=SLIDE_H)
        for lk in links:
            left = Emu(int(lk["x"] * PX_TO_EMU))
            top = Emu(int(lk["y"] * PX_TO_EMU))
            width = Emu(int(lk["w"] * PX_TO_EMU))
            height = Emu(int(lk["h"] * PX_TO_EMU))
            box = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
            box.fill.background()
            box.line.fill.background()
            box.click_action.hyperlink.address = lk["href"]
            # remove placeholder text
            box.text_frame.text = ""
            print(f"  link {lk['href']} @ ({lk['x']:.0f},{lk['y']:.0f})")

    prs.save(OUT)
    print(f"wrote {OUT} ({OUT.stat().st_size} bytes)")


if __name__ == "__main__":
    items = collect()
    build(items)
