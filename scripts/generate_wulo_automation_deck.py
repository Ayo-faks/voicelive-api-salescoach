from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
ASSETS = DOCS / "assets"
OUT_PPTX = DOCS / "wulo-automation-impact-presentation.pptx"
OUT_MD = DOCS / "wulo-automation-impact-presentation-outline.md"


class C:
    bg = RGBColor(0, 0, 0)
    ink = RGBColor(255, 255, 255)
    muted = RGBColor(176, 184, 196)
    line = RGBColor(55, 65, 81)
    white = RGBColor(255, 255, 255)
    card = RGBColor(20, 22, 28)
    card_alt = RGBColor(28, 31, 38)
    teal = RGBColor(45, 212, 191)
    blue = RGBColor(96, 165, 250)
    amber = RGBColor(251, 191, 36)
    green = RGBColor(74, 222, 128)
    red = RGBColor(248, 113, 113)
    teal_light = RGBColor(20, 35, 38)
    blue_light = RGBColor(20, 28, 45)
    amber_light = RGBColor(45, 33, 18)
    green_light = RGBColor(20, 38, 26)
    slate_light = RGBColor(28, 31, 38)


SLIDES = [
    {
        "section": "Opening",
        "title": "Wulo",
        "subtitle": "The speech therapy operating system",
        "tag": "Speech Therapy OS",
        "notes": [
            "Frame Wulo as a workflow automation project, not just an AI demo.",
            "The goal was to recover clinical time while improving practice quality at home.",
        ],
    },
    {
        "section": "Executive Summary",
        "title": "Wulo in one slide",
        "summary": [
            ("Problem", "Therapists lost up to 30 minutes reconstructing what happened between sessions."),
            ("Approach", "Design a speech therapy OS that captures practice, analyses it, and presents it back to the therapist."),
            ("Implementation", "Real-time voice practice, phoneme-aware assessment, session summaries, review dashboards."),
            ("Results", "Review time reduced from 30 minutes to 5 minutes; clinical phoneme performance improved by 22%."),
        ],
        "notes": [
            "This slide previews the judging criteria: clarity, problem solving, technical depth, and measured outcomes.",
        ],
    },
    {
        "section": "1. The Problem",
        "title": "A one-hour session was being consumed by missing context",
        "headline": "Therapists were forced to reconstruct the week before they could treat the child in front of them.",
        "bullets": [
            "Home practice was inconsistent, undocumented, and hard to verify.",
            "Parents wanted to help but often lacked time, structure, or clinical confidence.",
            "Children needed repeated practice, but not every home had someone available to coach it.",
            "The first 30 minutes of a 60-minute session could disappear into catch-up and guesswork.",
        ],
        "notes": [
            "The waste was not a single UI problem. It was a broken information loop between therapist, parent, and child.",
        ],
    },
    {
        "section": "Problem Quantified",
        "title": "Where the inefficiency came from",
        "chain": [
            ("Therapist sets targets", "Good clinical plan"),
            ("Home practice happens", "Low structure"),
            ("Evidence is missing", "No reliable trace"),
            ("Next session starts", "30 min reconstruction"),
        ],
        "notes": [
            "The opportunity was to convert unstructured home activity into useful clinical evidence.",
        ],
    },
    {
        "section": "2. The Approach",
        "title": "I designed Wulo as an operating system for therapist-led care",
        "headline": "The question was not what can AI do? It was what should the OS run on behalf of the therapist, and what must the therapist still own?",
        "bullets": [
            "Observed the workflow around session prep, home practice, review, and next-step planning.",
            "Separated clinical judgement from repeatable evidence-capture tasks.",
            "Prioritized automations with measurable time savings and low safety risk.",
            "Designed human-in-the-loop checkpoints so generated outputs supported, rather than replaced, therapists.",
        ],
        "notes": [
            "This shows decision making: automate the repetitive work, preserve clinical control.",
        ],
    },
    {
        "section": "Automation Choices",
        "title": "What the OS runs for the therapist",
        "matrix": [
            ("Guided home practice", "Give the child structured repetitions without needing a parent to lead every turn."),
            ("Transcript capture", "Preserve what happened instead of relying on memory or parent recall."),
            ("Pronunciation scoring", "Turn audio into word and phoneme-level signals therapists can inspect."),
            ("Session summarisation", "Compress raw interaction data into review-ready clinical context."),
            ("Next-session planning", "Use saved evidence to draft a plan for therapist approval."),
        ],
        "notes": [
            "Each automation maps to a bottleneck in the original workflow.",
        ],
    },
    {
        "section": "3. Implementation",
        "title": "OS architecture: realtime practice plus review intelligence",
        "image": "architecture.png",
        "bullets": [
            "React and TypeScript frontend for child practice and therapist dashboard workflows.",
            "Python Flask backend with WebSocket proxy for real-time audio and avatar sessions.",
            "Azure Voice Live for guided conversation and avatar delivery.",
            "Azure Speech for pronunciation assessment; Azure OpenAI for structured analysis and planning.",
            "Persistence layer for session history, child memory, recommendations, and audit-friendly review.",
        ],
        "notes": [
            "Explain the system as two loops: real-time practice and asynchronous therapist review.",
        ],
    },
    {
        "section": "Implementation Detail",
        "title": "Realtime practice loop",
        "image": "preview.png",
        "headline": "The child gets a supportive practice partner; the therapist gets structured evidence back.",
        "steps": [
            "Audio capture streams from browser to backend over WebSocket.",
            "Voice session returns assistant audio, avatar output, and transcriptions.",
            "Session completion triggers analysis of transcript, target words, and pronunciation data.",
            "Results are saved for therapist review before the next appointment.",
        ],
        "notes": [
            "This converts home practice from an invisible event into a replayable, analysable artifact.",
        ],
    },
    {
        "section": "Implementation Detail",
        "title": "Making speech AI work for clinical phonemes",
        "headline": "Off-the-shelf speech models were not enough for therapy-specific pronunciation targets.",
        "bullets": [
            "Built a target-sound layer around words, phonemes, repetitions, and child-appropriate prompts.",
            "Used Azure Speech pronunciation assessment for word-level accuracy, fluency, and completeness signals.",
            "Added phoneme-aware content and lexicon support so clinical targets were pronounced and evaluated more reliably.",
            "Compared baseline model behaviour against therapy-specific prompts and target lists on clinical phoneme examples.",
        ],
        "metric": ("+22%", "Improvement over the off-the-shelf speech model on clinical phoneme pronunciation tasks"),
        "notes": [
            "The 22 percent result is the technical credibility anchor. Emphasize iteration: baseline, error analysis, targeted adaptation, retest.",
        ],
    },
    {
        "section": "Implementation Detail",
        "title": "Therapist review dashboard",
        "image": "analysis.png",
        "headline": "Raw practice becomes decision-ready evidence.",
        "bullets": [
            "Per-session transcripts and assessment results are stored for review.",
            "Pronunciation scores surface word-level strengths and failure cases.",
            "AI-generated notes summarise effort, clarity, retries, and suggested practice focus.",
            "Therapist feedback stays authoritative and can be added after review.",
        ],
        "notes": [
            "This is where the time saving shows up: therapists review a structured record instead of reconstructing from scratch.",
        ],
    },
    {
        "section": "4. Results",
        "title": "Measured impact",
        "metrics": [
            ("30 min", "Before", "Manual reconstruction at start of a 60-minute therapy session", C.red),
            ("5 min", "After", "Review using Wulo session evidence and summaries", C.green),
            ("83%", "Faster review", "25 minutes returned to therapy time", C.blue),
            ("+22%", "Speech quality", "Improvement on clinical phoneme pronunciation tasks", C.teal),
        ],
        "notes": [
            "State the before and after clearly. 30 to 5 minutes is an 83 percent reduction and returns 25 minutes to direct therapy.",
        ],
    },
    {
        "section": "Why It Matters",
        "title": "The OS changed the economics of care delivery",
        "headline": "The value was not just speed. It was better allocation of scarce clinical attention.",
        "bullets": [
            "Performance: therapists start with evidence, not uncertainty.",
            "Efficiency: 25 minutes per session moves from reconstruction to intervention.",
            "Quality: phoneme-specific feedback improves the signal therapists use to adjust plans.",
            "Scalability: parents no longer need to personally lead every practice repetition.",
            "Safety: AI supports preparation and review while therapists retain clinical judgement.",
        ],
        "notes": [
            "Tie the outcome to performance, efficiency, cost, and quality as requested.",
        ],
    },
    {
        "section": "Closing",
        "title": "The engineering lesson",
        "headline": "A good speech therapy OS respects the expert workflow it runs on.",
        "bullets": [
            "I did not build an AI therapist. I built the OS that surrounds the therapist.",
            "It captures practice, analyses speech, summarises evidence, and proposes next steps.",
            "The therapist remains the decision maker, but now enters the session with better data and more time.",
        ],
        "closing": "Wulo is the speech therapy OS that turns between-session practice into measurable clinical progress.",
        "notes": [
            "End with the senior-engineer framing: the best technical decision was preserving the human control boundary.",
        ],
    },
]


def rgb(value: RGBColor) -> str:
    return f"#{value[0]:02x}{value[1]:02x}{value[2]:02x}"


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C.bg
    return slide


def add_box(slide, x, y, w, h, fill=C.card, line=C.line, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(1)
    return box


def add_text(slide, text, x, y, w, h, size=22, color=C.ink, bold=False, align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_rich_line(slide, label, body, x, y, w, h, label_color=C.teal, body_color=C.ink, size=18):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{label}: "
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = label_color
    run = p.add_run()
    run.text = body
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.color.rgb = body_color
    return tb


def add_header(slide, section, title, slide_num):
    add_text(slide, section.upper(), 0.62, 0.36, 3.0, 0.25, size=9, color=C.teal, bold=True)
    add_text(slide, title, 0.62, 0.67, 8.9, 0.7, size=29, bold=True)
    add_text(slide, f"{slide_num:02d}", 12.1, 0.44, 0.55, 0.25, size=9, color=C.muted, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.34), Inches(12.1), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = C.line
    line.line.color.rgb = C.line


def add_bullet_list(slide, bullets, x, y, w, h, size=17, color=C.ink, bullet_color=C.teal, gap=0.52):
    current_y = y
    for bullet in bullets:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(current_y + 0.09), Inches(0.13), Inches(0.13))
        dot.fill.solid()
        dot.fill.fore_color.rgb = bullet_color
        dot.line.color.rgb = bullet_color
        add_text(slide, bullet, x + 0.28, current_y, w - 0.28, 0.45, size=size, color=color)
        current_y += gap


def add_image_fit(slide, image_name, x, y, w, h):
    path = ASSETS / image_name
    if not path.exists():
        add_box(slide, x, y, w, h, fill=C.slate_light)
        add_text(slide, f"Missing image: {image_name}", x + 0.2, y + 0.2, w - 0.4, 0.4, size=14, color=C.muted)
        return
    with Image.open(path) as img:
        img_w, img_h = img.size
    ratio = min(w / img_w, h / img_h)
    pic_w = img_w * ratio
    pic_h = img_h * ratio
    pic_x = x + (w - pic_w) / 2
    pic_y = y + (h - pic_h) / 2
    add_box(slide, x, y, w, h, fill=C.card)
    slide.shapes.add_picture(str(path), Inches(pic_x), Inches(pic_y), Inches(pic_w), Inches(pic_h))


def slide_title(prs, item, idx):
    slide = blank_slide(prs)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(2.15))
    band.fill.solid()
    band.fill.fore_color.rgb = C.teal_light
    band.line.color.rgb = C.teal_light
    add_text(slide, item["tag"].upper(), 0.75, 0.78, 5.0, 0.3, size=11, color=C.teal, bold=True)
    add_text(slide, item["title"], 0.75, 1.45, 7.6, 0.95, size=56, bold=True, color=C.white)
    add_text(slide, item["subtitle"], 0.8, 2.55, 7.6, 0.75, size=24, color=C.muted)
    add_box(slide, 8.75, 1.65, 3.75, 3.65, fill=C.card_alt, line=C.line)
    add_text(slide, "30 -> 5 min", 9.05, 2.25, 3.1, 0.6, size=32, color=C.teal, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "home-practice review", 9.05, 2.95, 3.1, 0.35, size=15, color=C.muted, align=PP_ALIGN.CENTER)
    add_text(slide, "+22%", 9.05, 3.75, 3.1, 0.6, size=38, color=C.green, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "clinical phoneme performance", 9.05, 4.45, 3.1, 0.35, size=15, color=C.muted, align=PP_ALIGN.CENTER)
    add_text(slide, "Built with React, Python, realtime speech, Azure Speech, Azure Voice Live, and Azure OpenAI", 0.8, 6.6, 11.8, 0.35, size=14, color=C.muted)
    return slide


def slide_summary(prs, item, idx):
    slide = blank_slide(prs)
    add_header(slide, item["section"], item["title"], idx)
    colors = [C.teal_light, C.blue_light, C.amber_light, C.green_light]
    accent_colors = [C.teal, C.blue, C.amber, C.green]
    for i, (label, body) in enumerate(item["summary"]):
        x = 0.72 + (i % 2) * 6.15
        y = 1.82 + (i // 2) * 2.25
        add_box(slide, x, y, 5.55, 1.75, fill=colors[i], line=colors[i])
        add_text(slide, label, x + 0.28, y + 0.22, 2.1, 0.35, size=16, color=accent_colors[i], bold=True)
        add_text(slide, body, x + 0.28, y + 0.72, 4.95, 0.8, size=18, color=C.ink, bold=True)
    return slide


def slide_problem(prs, item, idx):
    slide = blank_slide(prs)
    add_header(slide, item["section"], item["title"], idx)
    add_text(slide, item["headline"], 0.72, 1.75, 6.2, 0.95, size=24, color=C.ink, bold=True)
    add_bullet_list(slide, item["bullets"], 0.78, 3.05, 6.25, 2.5, size=16)
    add_box(slide, 8.05, 1.85, 3.85, 3.75, fill=C.card_alt)
    add_text(slide, "The hidden cost", 8.4, 2.2, 3.15, 0.35, size=16, color=C.muted, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "50%", 8.55, 2.75, 2.85, 0.95, size=52, color=C.red, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "of a one-hour session could be spent reconstructing home practice", 8.55, 3.78, 2.85, 1.0, size=18, color=C.ink, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "30 minutes lost before intervention could properly begin", 8.55, 4.93, 2.85, 0.4, size=13, color=C.muted, align=PP_ALIGN.CENTER)
    return slide


def slide_chain(prs, item, idx):
    slide = blank_slide(prs)
    add_header(slide, item["section"], item["title"], idx)
    start_x = 0.65
    y = 2.3
    card_w = 2.75
    colors = [C.teal_light, C.amber_light, C.card_alt, C.blue_light]
    for i, (title, sub) in enumerate(item["chain"]):
        x = start_x + i * 3.12
        fill = colors[i]
        add_box(slide, x, y, card_w, 1.7, fill=fill, line=fill)
        add_text(slide, str(i + 1), x + 0.18, y + 0.18, 0.32, 0.3, size=13, color=C.muted, bold=True)
        add_text(slide, title, x + 0.35, y + 0.48, 2.05, 0.55, size=17, color=C.ink, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, sub, x + 0.35, y + 1.14, 2.05, 0.3, size=13, color=C.muted, align=PP_ALIGN.CENTER)
        if i < len(item["chain"]) - 1:
            add_text(slide, ">", x + 2.82, y + 0.64, 0.35, 0.35, size=26, color=C.muted, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 1.18, 5.15, 10.95, 0.86, fill=C.card_alt)
    add_text(slide, "Automation target", 1.45, 5.35, 1.9, 0.3, size=13, color=C.teal, bold=True)
    add_text(slide, "Capture the practice evidence automatically, then compress it into something a therapist can review in minutes.", 3.15, 5.32, 8.5, 0.35, size=18, color=C.ink, bold=True)
    return slide


def slide_text(prs, item, idx):
    slide = blank_slide(prs)
    add_header(slide, item["section"], item["title"], idx)
    add_text(slide, item["headline"], 0.72, 1.75, 11.1, 0.75, size=24, color=C.ink, bold=True)
    add_bullet_list(slide, item["bullets"], 0.82, 2.9, 10.7, 2.6, size=17)
    return slide


def slide_matrix(prs, item, idx):
    slide = blank_slide(prs)
    add_header(slide, item["section"], item["title"], idx)
    fills = [C.teal_light, C.blue_light, C.amber_light, C.green_light, C.slate_light]
    for i, (label, body) in enumerate(item["matrix"]):
        y = 1.75 + i * 0.92
        add_box(slide, 0.75, y, 11.78, 0.68, fill=fills[i], line=fills[i])
        add_text(slide, label, 1.05, y + 0.17, 2.75, 0.25, size=15, color=C.ink, bold=True)
        add_text(slide, body, 3.9, y + 0.17, 8.05, 0.25, size=14, color=C.muted)
    return slide


def slide_architecture(prs, item, idx):
    slide = blank_slide(prs)
    add_header(slide, item["section"], item["title"], idx)
    add_image_fit(slide, item["image"], 0.75, 1.68, 6.45, 4.75)
    add_bullet_list(slide, item["bullets"], 7.55, 1.85, 4.95, 3.9, size=14, gap=0.68, bullet_color=C.blue)
    return slide


def slide_realtime(prs, item, idx):
    slide = blank_slide(prs)
    add_header(slide, item["section"], item["title"], idx)
    add_image_fit(slide, item["image"], 0.75, 1.7, 5.8, 4.85)
    add_text(slide, item["headline"], 6.95, 1.78, 5.2, 0.75, size=22, color=C.ink, bold=True)
    for i, step in enumerate(item["steps"]):
        y = 2.85 + i * 0.82
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.0), Inches(y), Inches(0.36), Inches(0.36))
        circle.fill.solid()
        circle.fill.fore_color.rgb = C.teal
        circle.line.color.rgb = C.teal
        add_text(slide, str(i + 1), 7.0, y + 0.05, 0.36, 0.15, size=9, color=C.white, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, step, 7.55, y - 0.02, 4.8, 0.45, size=15, color=C.muted)
    return slide


def slide_phoneme(prs, item, idx):
    slide = blank_slide(prs)
    add_header(slide, item["section"], item["title"], idx)
    add_text(slide, item["headline"], 0.75, 1.72, 7.1, 0.7, size=22, color=C.ink, bold=True)
    add_bullet_list(slide, item["bullets"], 0.82, 2.75, 6.95, 2.6, size=15, gap=0.61)
    add_box(slide, 8.35, 1.98, 3.65, 3.65, fill=C.card_alt, line=C.line)
    metric, label = item["metric"]
    add_text(slide, metric, 8.75, 2.58, 2.85, 0.75, size=52, color=C.green, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, 8.75, 3.55, 2.85, 1.0, size=16, color=C.white, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Measured by comparing baseline speech behaviour with therapy-specific phoneme prompts and target lists.", 8.75, 4.78, 2.85, 0.48, size=11, color=C.muted, align=PP_ALIGN.CENTER)
    return slide


def slide_dashboard(prs, item, idx):
    slide = blank_slide(prs)
    add_header(slide, item["section"], item["title"], idx)
    add_text(slide, item["headline"], 0.75, 1.68, 5.3, 0.55, size=22, color=C.ink, bold=True)
    add_bullet_list(slide, item["bullets"], 0.85, 2.55, 5.25, 2.65, size=15, gap=0.62, bullet_color=C.blue)
    add_image_fit(slide, item["image"], 6.35, 1.68, 6.15, 4.85)
    return slide


def slide_metrics(prs, item, idx):
    slide = blank_slide(prs)
    add_header(slide, item["section"], item["title"], idx)
    for i, (big, label, body, color) in enumerate(item["metrics"]):
        x = 0.78 + (i % 2) * 6.1
        y = 1.75 + (i // 2) * 2.15
        add_box(slide, x, y, 5.55, 1.65, fill=C.card_alt)
        add_text(slide, big, x + 0.3, y + 0.25, 1.85, 0.55, size=31, color=color, bold=True)
        add_text(slide, label, x + 2.2, y + 0.27, 2.8, 0.3, size=16, color=C.ink, bold=True)
        add_text(slide, body, x + 2.2, y + 0.78, 2.8, 0.45, size=13, color=C.muted)
    add_box(slide, 1.22, 6.0, 10.85, 0.58, fill=C.teal_light, line=C.teal_light)
    add_text(slide, "Bottom line: Wulo returned 25 minutes of a 60-minute therapy session to active clinical work.", 1.45, 6.16, 10.35, 0.24, size=16, color=C.teal, bold=True, align=PP_ALIGN.CENTER)
    return slide


def slide_closing(prs, item, idx):
    slide = slide_text(prs, item, idx)
    add_box(slide, 0.85, 6.02, 11.55, 0.65, fill=C.teal, line=C.teal)
    add_text(slide, item["closing"], 1.15, 6.2, 10.95, 0.25, size=17, color=C.bg, bold=True, align=PP_ALIGN.CENTER)
    return slide


def build_outline():
    lines = ["# Wulo Automation Impact Presentation", "", "Generated source outline for presenter rehearsal and quick edits.", ""]
    for i, item in enumerate(SLIDES, start=1):
        lines.append(f"## {i}. {item['title']}")
        lines.append(f"Section: {item['section']}")
        if "subtitle" in item:
            lines.append(f"Subtitle: {item['subtitle']}")
        if "headline" in item:
            lines.append(f"Headline: {item['headline']}")
        for key in ("bullets", "steps"):
            if key in item:
                lines.append("")
                for value in item[key]:
                    lines.append(f"- {value}")
        if "summary" in item:
            lines.append("")
            for label, body in item["summary"]:
                lines.append(f"- {label}: {body}")
        if "matrix" in item:
            lines.append("")
            for label, body in item["matrix"]:
                lines.append(f"- {label}: {body}")
        if "chain" in item:
            lines.append("")
            for label, body in item["chain"]:
                lines.append(f"- {label}: {body}")
        if "metrics" in item:
            lines.append("")
            for big, label, body, _ in item["metrics"]:
                lines.append(f"- {big} {label}: {body}")
        if "metric" in item:
            big, body = item["metric"]
            lines.append(f"- {big}: {body}")
        if "closing" in item:
            lines.append(f"Closing: {item['closing']}")
        if "notes" in item:
            lines.append("")
            lines.append("Presenter notes:")
            for note in item["notes"]:
                lines.append(f"- {note}")
        lines.append("")
    return "\n".join(lines)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_title,
        slide_summary,
        slide_problem,
        slide_chain,
        slide_text,
        slide_matrix,
        slide_architecture,
        slide_realtime,
        slide_phoneme,
        slide_dashboard,
        slide_metrics,
        slide_text,
        slide_closing,
    ]

    for idx, (builder, item) in enumerate(zip(builders, SLIDES), start=1):
        builder(prs, item, idx)

    prs.save(OUT_PPTX)
    OUT_MD.write_text(build_outline(), encoding="utf-8")
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {OUT_MD}")


if __name__ == "__main__":
    main()